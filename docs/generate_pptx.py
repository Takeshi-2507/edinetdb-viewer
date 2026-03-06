"""
SnowWillow Terminal システム構成図 PowerPoint 生成スクリプト
"""
import requests
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

ICONS_DIR = Path(__file__).parent / "icons"
ICONS_DIR.mkdir(exist_ok=True)

# --------------- Logo Downloads ---------------
LOGOS = {
    "react": "https://upload.wikimedia.org/wikipedia/commons/thumb/a/a7/React-icon.svg/512px-React-icon.svg.png",
    "python": "https://upload.wikimedia.org/wikipedia/commons/thumb/c/c3/Python-logo-notext.svg/480px-Python-logo-notext.svg.png",
    "sqlite": "https://upload.wikimedia.org/wikipedia/commons/thumb/9/97/Sqlite-square-icon.svg/480px-Sqlite-square-icon.svg.png",
    "github": "https://github.githubassets.com/images/modules/logos_page/GitHub-Mark.png",
    "docker": "https://cdn4.iconfinder.com/data/icons/logos-and-brands/512/97_Docker_logo_logos-512.png",
    "vite": "https://vitejs.dev/logo.svg",
}

def download_logos():
    for name, url in LOGOS.items():
        ext = "png" if url.endswith(".png") else "svg"
        path = ICONS_DIR / f"{name}.png"
        if path.exists():
            continue
        try:
            r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            if r.status_code == 200:
                path.write_bytes(r.content)
                print(f"  Downloaded {name}")
            else:
                print(f"  Failed {name}: {r.status_code}")
        except Exception as e:
            print(f"  Error {name}: {e}")

# --------------- Color Scheme ---------------
BG_DARK = RGBColor(0x08, 0x08, 0x0C)
BG_SURFACE = RGBColor(0x12, 0x12, 0x1A)
BG_SURFACE2 = RGBColor(0x1C, 0x1C, 0x28)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6A, 0x00)
WHITE = RGBColor(0xF0, 0xF0, 0xF5)
GRAY = RGBColor(0x6E, 0x6E, 0x82)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_text_in_shape(shape, text, font_size=11, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return shape

def add_arrow(slide, start_x, start_y, end_x, end_y, color=GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.begin_x = start_x
    connector.begin_y = start_y
    connector.end_x = end_x
    connector.end_y = end_y
    return connector

def add_icon(slide, name, left, top, size=Inches(0.4)):
    path = ICONS_DIR / f"{name}.png"
    if path.exists() and path.stat().st_size > 0:
        try:
            slide.shapes.add_picture(str(path), left, top, size, size)
            return True
        except Exception:
            pass
    return False

# --------------- Slide 1: Local Version ---------------
def build_local_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Title
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9), Inches(0.5),
                 "SnowWillow Terminal — ローカル版 システム構成図",
                 font_size=22, color=CYAN, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.4), Inches(0.6), Inches(9), Inches(0.3),
                 "PC (Windows 11) 上で完結  |  株価リアルタイム  |  サイバーテーマ",
                 font_size=11, color=GRAY)

    # ===== Browser / Frontend =====
    browser_box = add_rounded_box(slide,
        Inches(0.3), Inches(1.1), Inches(4.2), Inches(2.4),
        BG_SURFACE, CYAN, Pt(2))
    add_text_in_shape(browser_box, "", font_size=9)

    add_text_box(slide, Inches(0.5), Inches(1.15), Inches(3), Inches(0.3),
                 "ブラウザ  localhost:5173", font_size=13, color=CYAN, bold=True)

    # Frontend components
    fe_box = add_rounded_box(slide,
        Inches(0.5), Inches(1.55), Inches(3.8), Inches(1.8),
        BG_SURFACE2, RGBColor(0x2A, 0x2A, 0x3A))

    add_icon(slide, "react", Inches(0.6), Inches(1.6), Inches(0.35))
    add_text_box(slide, Inches(1.0), Inches(1.58), Inches(2), Inches(0.25),
                 "React 18 + Vite", font_size=11, color=WHITE, bold=True)

    components = [
        ("Screener", "V/Qスコア統合", CYAN),
        ("CompanyDetail", "財務チャート+IR", WHITE),
        ("Dashboard", "ステータス+業種分析", WHITE),
        ("DemoTrade", "模擬売買+ポートフォリオ", WHITE),
    ]
    for i, (name, desc, col) in enumerate(components):
        y = Inches(1.95) + Inches(i * 0.32)
        add_text_box(slide, Inches(0.65), y, Inches(3.5), Inches(0.28),
                     f"● {name}  {desc}", font_size=9, color=col)

    # Theme badge
    theme_box = add_rounded_box(slide,
        Inches(2.8), Inches(1.15), Inches(1.5), Inches(0.28),
        RGBColor(0x00, 0x20, 0x30), CYAN, Pt(1))
    add_text_in_shape(theme_box, "🌍 サイバーテーマ", font_size=8, color=CYAN)

    # ===== Arrow: Frontend -> Backend =====
    add_text_box(slide, Inches(1.5), Inches(3.55), Inches(2), Inches(0.25),
                 "▼  /api/*  (Vite Proxy)", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ===== Backend =====
    be_box = add_rounded_box(slide,
        Inches(0.3), Inches(3.85), Inches(4.2), Inches(2.8),
        BG_SURFACE, PURPLE, Pt(2))

    add_icon(slide, "python", Inches(0.5), Inches(3.9), Inches(0.35))
    add_text_box(slide, Inches(0.9), Inches(3.9), Inches(3), Inches(0.3),
                 "FastAPI + Uvicorn  :8001", font_size=13, color=PURPLE, bold=True)

    scoring_items = [
        ("Value Score (竹原式)", "PER+PBR+ROE+営業利益率+現金+FCF → 100点", ORANGE),
        ("Quality Score", "営業利益率+ROE+CF質 → 100点", GREEN),
        ("統合スコア", "Value×0.6 + Quality×0.4", CYAN),
        ("偽成長検知", "4フラグ → Quality減点", YELLOW),
        ("業種除外フィルタ", "複数業種チェックボックス", WHITE),
        ("CN-PER", "PER×(1-ネットキャッシュ比率)", WHITE),
        ("売りアラート", "EPS×15超えで通知", RED),
    ]
    for i, (name, desc, col) in enumerate(scoring_items):
        y = Inches(4.35) + Inches(i * 0.3)
        add_text_box(slide, Inches(0.55), y, Inches(1.8), Inches(0.25),
                     f"▸ {name}", font_size=9, color=col, bold=True)
        add_text_box(slide, Inches(2.35), y, Inches(2.1), Inches(0.25),
                     desc, font_size=8, color=GRAY)

    # ===== Right Side: Data Sources =====

    # SQLite
    db_box = add_rounded_box(slide,
        Inches(5.2), Inches(1.1), Inches(2.2), Inches(2.0),
        BG_SURFACE, GREEN, Pt(1.5))
    add_icon(slide, "sqlite", Inches(5.35), Inches(1.18), Inches(0.35))
    add_text_box(slide, Inches(5.75), Inches(1.18), Inches(1.5), Inches(0.3),
                 "SQLite", font_size=13, color=GREEN, bold=True)

    db_items = [
        "edinet.db",
        "3,848社 / 財務10年分",
        "companies テーブル",
        "financials テーブル",
        "ratios テーブル",
        "analysis テーブル",
    ]
    for i, item in enumerate(db_items):
        add_text_box(slide, Inches(5.4), Inches(1.55 + i * 0.23), Inches(1.9), Inches(0.22),
                     item, font_size=8, color=WHITE if i < 2 else GRAY)

    # yfinance
    yf_box = add_rounded_box(slide,
        Inches(5.2), Inches(3.3), Inches(2.2), Inches(1.6),
        BG_SURFACE, YELLOW, Pt(1.5))
    add_text_box(slide, Inches(5.4), Inches(3.38), Inches(1.8), Inches(0.3),
                 "yfinance", font_size=13, color=YELLOW, bold=True)

    yf_items = [
        "リアルタイム株価",
        "バランスシート (BS)",
        "TTL: 5分 / BS: 30日",
        "並列取得 (5workers)",
        "→ Yahoo Finance API",
    ]
    for i, item in enumerate(yf_items):
        add_text_box(slide, Inches(5.4), Inches(3.72 + i * 0.22), Inches(1.9), Inches(0.2),
                     item, font_size=8, color=WHITE if i < 2 else GRAY)

    # Cache files
    cache_box = add_rounded_box(slide,
        Inches(5.2), Inches(5.1), Inches(2.2), Inches(0.9),
        BG_SURFACE, GRAY, Pt(1))
    add_text_box(slide, Inches(5.4), Inches(5.15), Inches(1.8), Inches(0.25),
                 "キャッシュファイル", font_size=10, color=GRAY, bold=True)
    add_text_box(slide, Inches(5.4), Inches(5.4), Inches(1.8), Inches(0.2),
                 "price_cache.json (5分)", font_size=8, color=GRAY)
    add_text_box(slide, Inches(5.4), Inches(5.6), Inches(1.8), Inches(0.2),
                 "bs_cache.json (30日)", font_size=8, color=GRAY)

    # ===== Right side: External =====
    # EDINET
    ext_box = add_rounded_box(slide,
        Inches(8.0), Inches(1.1), Inches(1.7), Inches(1.2),
        BG_SURFACE2, RGBColor(0x40, 0x40, 0x60), Pt(1))
    add_text_box(slide, Inches(8.1), Inches(1.18), Inches(1.5), Inches(0.25),
                 "EDINET DB API", font_size=10, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.1), Inches(1.45), Inches(1.5), Inches(0.2),
                 "edinetdb.jp", font_size=8, color=GRAY)
    add_text_box(slide, Inches(8.1), Inches(1.65), Inches(1.5), Inches(0.2),
                 "毎日 JST 3:00 同期", font_size=8, color=CYAN)
    add_text_box(slide, Inches(8.1), Inches(1.85), Inches(1.5), Inches(0.2),
                 "(GitHub Actions)", font_size=8, color=GRAY)

    # Yahoo Finance
    yh_box = add_rounded_box(slide,
        Inches(8.0), Inches(2.6), Inches(1.7), Inches(0.9),
        BG_SURFACE2, RGBColor(0x40, 0x40, 0x60), Pt(1))
    add_text_box(slide, Inches(8.1), Inches(2.68), Inches(1.5), Inches(0.25),
                 "Yahoo Finance", font_size=10, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.1), Inches(2.95), Inches(1.5), Inches(0.2),
                 "株価 + 財務データ", font_size=8, color=GRAY)
    add_text_box(slide, Inches(8.1), Inches(3.15), Inches(1.5), Inches(0.2),
                 "HTTPS リアルタイム", font_size=8, color=YELLOW)

    # Connection arrows (text-based)
    add_text_box(slide, Inches(4.55), Inches(2.0), Inches(0.6), Inches(0.2),
                 "◀──", font_size=10, color=GREEN)
    add_text_box(slide, Inches(4.55), Inches(3.9), Inches(0.6), Inches(0.2),
                 "◀──", font_size=10, color=YELLOW)
    add_text_box(slide, Inches(7.45), Inches(1.5), Inches(0.5), Inches(0.2),
                 "──▶", font_size=10, color=GRAY)
    add_text_box(slide, Inches(7.45), Inches(2.9), Inches(0.5), Inches(0.2),
                 "──▶", font_size=10, color=YELLOW)

    # ===== Footer: Phase roadmap =====
    add_rounded_box(slide, Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.55),
                    BG_SURFACE2, RGBColor(0x2A, 0x2A, 0x3A), Pt(1))

    phases = [
        ("A Value ✅", ORANGE),
        ("B Quality ✅", GREEN),
        ("C Momentum 🔜", GRAY),
        ("D Event 🔜", GRAY),
        ("E AI定性 🔜", GRAY),
    ]
    for i, (label, col) in enumerate(phases):
        add_text_box(slide, Inches(0.5 + i * 1.85), Inches(6.9), Inches(1.7), Inches(0.25),
                     label, font_size=10, color=col, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.2),
                 "SnowWillow Terminal ロードマップ:  Phase 1 完了  →  Phase 2 (株価時系列/テクニカル) は次回",
                 font_size=8, color=GRAY)


# --------------- Slide 2: Online Version ---------------
def build_online_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9), Inches(0.5),
                 "SnowWillow Terminal — オンライン版 システム構成図",
                 font_size=22, color=ORANGE, bold=True)
    add_text_box(slide, Inches(0.4), Inches(0.6), Inches(9), Inches(0.3),
                 "Render.com (Free tier)  |  Docker  |  GitHub Actions 自動デプロイ  |  Asiimovテーマ",
                 font_size=11, color=GRAY)

    # ===== User =====
    user_box = add_rounded_box(slide,
        Inches(0.3), Inches(1.3), Inches(1.8), Inches(1.2),
        BG_SURFACE, GRAY, Pt(1.5))
    add_text_box(slide, Inches(0.4), Inches(1.38), Inches(1.6), Inches(0.3),
                 "👤 ユーザー", font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.4), Inches(1.7), Inches(1.6), Inches(0.2),
                 "ブラウザ (HTTPS)", font_size=9, color=GRAY)
    add_text_box(slide, Inches(0.4), Inches(1.95), Inches(1.6), Inches(0.2),
                 "テーマ: Asiimov (orange)", font_size=8, color=ORANGE)
    add_text_box(slide, Inches(0.4), Inches(2.15), Inches(1.6), Inches(0.2),
                 "株価: デフォルトOFF", font_size=8, color=GRAY)

    # Arrow user -> Render
    add_text_box(slide, Inches(2.15), Inches(1.7), Inches(0.5), Inches(0.2),
                 "──▶", font_size=12, color=GRAY)
    add_text_box(slide, Inches(2.15), Inches(1.9), Inches(0.6), Inches(0.2),
                 "HTTPS", font_size=8, color=GRAY)

    # ===== Render.com =====
    render_box = add_rounded_box(slide,
        Inches(2.7), Inches(1.1), Inches(4.5), Inches(3.8),
        BG_SURFACE, ORANGE, Pt(2))

    add_text_box(slide, Inches(2.9), Inches(1.18), Inches(3), Inches(0.3),
                 "Render.com (Free tier)", font_size=14, color=ORANGE, bold=True)

    # Docker container
    docker_box = add_rounded_box(slide,
        Inches(2.9), Inches(1.6), Inches(4.1), Inches(3.1),
        BG_SURFACE2, RGBColor(0x00, 0x7B, 0xFF), Pt(1.5))

    add_icon(slide, "docker", Inches(3.0), Inches(1.65), Inches(0.3))
    add_text_box(slide, Inches(3.35), Inches(1.65), Inches(2), Inches(0.25),
                 "Docker コンテナ", font_size=11, color=RGBColor(0x00, 0x7B, 0xFF), bold=True)
    add_text_box(slide, Inches(5.3), Inches(1.65), Inches(1.5), Inches(0.25),
                 "メモリ: 512MB", font_size=8, color=RED)

    # Uvicorn + FastAPI
    api_box = add_rounded_box(slide,
        Inches(3.1), Inches(2.0), Inches(3.7), Inches(1.2),
        RGBColor(0x15, 0x15, 0x22), PURPLE, Pt(1))
    add_icon(slide, "python", Inches(3.2), Inches(2.05), Inches(0.3))
    add_text_box(slide, Inches(3.55), Inches(2.05), Inches(2.5), Inches(0.25),
                 "Uvicorn + FastAPI  :8080", font_size=11, color=PURPLE, bold=True)

    api_items = [
        ("全スコアリング", "ローカルと同一コード"),
        ("業種除外/偽成長検知", "全フィルタ対応"),
        ("IR プロキシ", "EDINET API 中継"),
    ]
    for i, (name, desc) in enumerate(api_items):
        add_text_box(slide, Inches(3.2), Inches(2.38 + i * 0.25), Inches(1.6), Inches(0.22),
                     f"▸ {name}", font_size=8, color=WHITE, bold=True)
        add_text_box(slide, Inches(4.85), Inches(2.38 + i * 0.25), Inches(1.8), Inches(0.22),
                     desc, font_size=8, color=GRAY)

    # Static files
    static_box = add_rounded_box(slide,
        Inches(3.1), Inches(3.3), Inches(1.7), Inches(0.8),
        RGBColor(0x15, 0x15, 0x22), CYAN, Pt(1))
    add_icon(slide, "react", Inches(3.2), Inches(3.35), Inches(0.25))
    add_text_box(slide, Inches(3.5), Inches(3.35), Inches(1.2), Inches(0.25),
                 "React Build", font_size=10, color=CYAN, bold=True)
    add_text_box(slide, Inches(3.2), Inches(3.65), Inches(1.5), Inches(0.2),
                 "Static Files 配信", font_size=8, color=GRAY)
    add_text_box(slide, Inches(3.2), Inches(3.85), Inches(1.5), Inches(0.2),
                 "Asiimov テーマ", font_size=8, color=ORANGE)

    # SQLite (in container)
    db_box = add_rounded_box(slide,
        Inches(5.1), Inches(3.3), Inches(1.7), Inches(0.8),
        RGBColor(0x15, 0x15, 0x22), GREEN, Pt(1))
    add_icon(slide, "sqlite", Inches(5.2), Inches(3.35), Inches(0.25))
    add_text_box(slide, Inches(5.5), Inches(3.35), Inches(1.2), Inches(0.25),
                 "SQLite", font_size=10, color=GREEN, bold=True)
    add_text_box(slide, Inches(5.2), Inches(3.65), Inches(1.5), Inches(0.2),
                 "edinet.db (コンテナ内)", font_size=8, color=GRAY)
    add_text_box(slide, Inches(5.2), Inches(3.85), Inches(1.5), Inches(0.2),
                 "デプロイ時に更新", font_size=8, color=GRAY)

    # ===== GitHub =====
    gh_box = add_rounded_box(slide,
        Inches(2.7), Inches(5.3), Inches(2.0), Inches(1.8),
        BG_SURFACE, WHITE, Pt(1.5))
    add_icon(slide, "github", Inches(2.85), Inches(5.38), Inches(0.35))
    add_text_box(slide, Inches(3.25), Inches(5.38), Inches(1.3), Inches(0.3),
                 "GitHub", font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.85), Inches(5.75), Inches(1.7), Inches(0.2),
                 "main ブランチ", font_size=9, color=GRAY)
    add_text_box(slide, Inches(2.85), Inches(6.0), Inches(1.7), Inches(0.2),
                 "push → 自動デプロイ", font_size=8, color=ORANGE)
    add_text_box(slide, Inches(2.85), Inches(6.25), Inches(1.7), Inches(0.2),
                 "ソースコード管理", font_size=8, color=GRAY)
    add_text_box(slide, Inches(2.85), Inches(6.5), Inches(1.7), Inches(0.2),
                 "edinet.db も管理", font_size=8, color=GRAY)

    # Arrow GitHub -> Render
    add_text_box(slide, Inches(3.3), Inches(5.05), Inches(1.5), Inches(0.22),
                 "▲ auto deploy", font_size=8, color=ORANGE, alignment=PP_ALIGN.CENTER)

    # ===== GitHub Actions =====
    ga_box = add_rounded_box(slide,
        Inches(5.2), Inches(5.3), Inches(2.3), Inches(1.8),
        BG_SURFACE, CYAN, Pt(1.5))
    add_text_box(slide, Inches(5.35), Inches(5.38), Inches(2), Inches(0.3),
                 "GitHub Actions", font_size=12, color=CYAN, bold=True)
    add_text_box(slide, Inches(5.35), Inches(5.7), Inches(2), Inches(0.2),
                 "⏰ 毎日 JST 3:00", font_size=9, color=WHITE)

    ga_steps = [
        "1. collector.py 実行",
        "2. EDINET API → DB更新",
        "3. git commit + push",
        "4. → Render 自動デプロイ",
    ]
    for i, step in enumerate(ga_steps):
        add_text_box(slide, Inches(5.35), Inches(5.98 + i * 0.22), Inches(2), Inches(0.2),
                     step, font_size=8, color=GRAY if i < 3 else ORANGE)

    # Arrow Actions -> GitHub
    add_text_box(slide, Inches(4.75), Inches(6.2), Inches(0.5), Inches(0.2),
                 "◀──", font_size=10, color=CYAN)

    # ===== EDINET API =====
    edinet_box = add_rounded_box(slide,
        Inches(8.0), Inches(5.5), Inches(1.7), Inches(1.2),
        BG_SURFACE2, RGBColor(0x40, 0x40, 0x60), Pt(1))
    add_text_box(slide, Inches(8.1), Inches(5.58), Inches(1.5), Inches(0.25),
                 "EDINET DB API", font_size=10, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.1), Inches(5.85), Inches(1.5), Inches(0.2),
                 "edinetdb.jp", font_size=9, color=GRAY)
    add_text_box(slide, Inches(8.1), Inches(6.1), Inches(1.5), Inches(0.2),
                 "財務データ取得", font_size=8, color=GRAY)
    add_text_box(slide, Inches(8.1), Inches(6.3), Inches(1.5), Inches(0.2),
                 "HTTPS", font_size=8, color=CYAN)

    # Arrow Actions -> EDINET
    add_text_box(slide, Inches(7.55), Inches(6.0), Inches(0.5), Inches(0.2),
                 "──▶", font_size=10, color=GRAY)

    # ===== Comparison table =====
    # (bottom of slide not needed, already covered in separate slide)


# --------------- Slide 3: Comparison Table ---------------
def build_comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.4), Inches(0.3), Inches(9), Inches(0.5),
                 "ローカル版 vs オンライン版  比較", font_size=22, color=WHITE, bold=True)

    headers = ["", "ローカル版", "オンライン版"]
    rows = [
        ["テーマ", "サイバー (cyan/purple) + 地球", "Asiimov (orange/black)"],
        ["株価表示", "デフォルト ON", "デフォルト OFF (メモリ節約)"],
        ["フロントエンド", "Vite dev server (HMR)", "Static build (Uvicorn配信)"],
        ["バックエンド", "Uvicorn :8001", "Uvicorn :8080 (Docker)"],
        ["DB更新", "git pull で反映", "GitHub Actions → 自動デプロイ"],
        ["メモリ制限", "なし (PC依存)", "512MB (Free tier)"],
        ["スコアリング", "同一コード", "同一コード"],
        ["アクセス", "localhost のみ", "インターネット公開"],
        ["用途", "メイン利用 (全機能)", "外出先/共有用"],
    ]

    col_widths = [Inches(2.0), Inches(3.5), Inches(3.5)]
    start_x = Inches(0.5)
    start_y = Inches(1.1)
    row_h = Inches(0.42)

    # Header row
    x = start_x
    for i, (header, w) in enumerate(zip(headers, col_widths)):
        box = add_rounded_box(slide, x, start_y, w, row_h,
                              BG_SURFACE2 if i == 0 else (CYAN if i == 1 else ORANGE),
                              None)
        col = WHITE if i == 0 else RGBColor(0x08, 0x08, 0x0C)
        add_text_in_shape(box, header, font_size=12, color=col, bold=True)
        x += w + Inches(0.05)

    # Data rows
    for ri, row in enumerate(rows):
        y = start_y + (ri + 1) * (row_h + Inches(0.04))
        x = start_x
        bg = BG_SURFACE if ri % 2 == 0 else BG_SURFACE2
        for ci, (cell, w) in enumerate(zip(row, col_widths)):
            box = add_rounded_box(slide, x, y, w, row_h, bg, None)
            col = GRAY if ci == 0 else WHITE
            align = PP_ALIGN.LEFT
            tf = box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = align
            p = tf.paragraphs[0]
            p.text = f"  {cell}"
            p.font.size = Pt(10)
            p.font.color.rgb = col
            p.font.bold = (ci == 0)
            x += w + Inches(0.05)


# --------------- Main ---------------
def main():
    print("Downloading logos...")
    download_logos()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    build_local_slide(prs)
    build_online_slide(prs)
    build_comparison_slide(prs)

    out_path = Path(__file__).parent / "SnowWillow_System_Architecture.pptx"
    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")

if __name__ == "__main__":
    main()
